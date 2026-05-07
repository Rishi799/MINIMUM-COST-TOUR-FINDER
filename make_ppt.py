
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT1    = RGBColor(0x00, 0xC9, 0xFF)   # Cyan
ACCENT2    = RGBColor(0xF7, 0x97, 0x1E)   # Amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG    = RGBColor(0x16, 0x2D, 0x44)   # Slightly lighter navy

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # completely blank

# ─── Helpers ────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def fill_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(18), bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_gradient_bar(slide, l, t, w, h, color=ACCENT1):
    """Thin accent line."""
    add_rect(slide, l, t, w, h, fill_color=color)

def section_header(slide, title, subtitle=""):
    fill_bg(slide)
    # Top accent bar
    add_gradient_bar(slide, 0, 0, 13.33, 0.07, ACCENT1)
    # Bottom accent bar
    add_gradient_bar(slide, 0, 7.43, 13.33, 0.07, ACCENT2)
    # Card
    add_rect(slide, 1.2, 1.8, 10.9, 3.9, fill_color=CARD_BG, line_color=ACCENT1, line_width=Pt(1.5))
    add_textbox(slide, title, 1.5, 2.2, 10.3, 1.8,
                font_size=Pt(40), bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, subtitle, 1.5, 3.7, 10.3, 1.2,
                    font_size=Pt(20), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

def content_slide_header(slide, title):
    fill_bg(slide)
    add_gradient_bar(slide, 0, 0, 13.33, 0.07, ACCENT1)
    add_gradient_bar(slide, 0, 7.43, 13.33, 0.07, ACCENT2)
    # Title background
    add_rect(slide, 0, 0.07, 13.33, 0.93, fill_color=CARD_BG)
    add_textbox(slide, title, 0.3, 0.12, 12.7, 0.85,
                font_size=Pt(28), bold=True, color=ACCENT1, align=PP_ALIGN.LEFT)
    # Divider
    add_rect(slide, 0.3, 0.97, 12.73, 0.04, fill_color=ACCENT2)

def add_bullet_card(slide, l, t, w, h, title, bullets, title_color=ACCENT1):
    add_rect(slide, l, t, w, h, fill_color=CARD_BG, line_color=ACCENT1, line_width=Pt(1))
    add_textbox(slide, title, l+0.15, t+0.12, w-0.3, 0.45,
                font_size=Pt(16), bold=True, color=title_color)
    # Accent line under card title
    add_rect(slide, l+0.15, t+0.55, w-0.3, 0.03, fill_color=ACCENT2)
    tb = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.65), Inches(w-0.3), Inches(h-0.8))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {b}"
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT_GREY


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title Slide
# ════════════════════════════════════════════════════════════════════════════
s1 = add_slide()
fill_bg(s1)
add_gradient_bar(s1, 0, 0,     13.33, 0.12, ACCENT1)
add_gradient_bar(s1, 0, 7.38,  13.33, 0.12, ACCENT2)

# Big decorative circle
c = s1.shapes.add_shape(9, Inches(9.2), Inches(0.4), Inches(4.5), Inches(4.5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x00,0x4A,0x6E)
c.line.fill.background()

add_textbox(s1, "MINIMUM COST TOUR FINDER", 0.5, 1.5, 9.5, 1.1,
            font_size=Pt(38), bold=True, color=ACCENT1)
add_textbox(s1, "FOR MULTI-CITY TRAVEL NETWORKS", 0.5, 2.55, 9.5, 0.9,
            font_size=Pt(26), bold=True, color=WHITE)
add_rect(s1, 0.5, 3.55, 5.5, 0.06, fill_color=ACCENT2)

add_textbox(s1, "Design and Analysis of Algorithms  |  CD343AI", 0.5, 3.75, 9.5, 0.5,
            font_size=Pt(15), color=LIGHT_GREY)
add_textbox(s1, "Department of Artificial Intelligence and Machine Learning", 0.5, 4.2, 9.5, 0.5,
            font_size=Pt(13), italic=True, color=LIGHT_GREY)

add_textbox(s1, "Rhea Parthiban  •  Rishi Agarwal  •  Rangappagari John Niranjan",
            0.5, 5.2, 9.5, 0.5, font_size=Pt(15), bold=True, color=ACCENT2)
add_textbox(s1, "1RV24CI097  •  1RV24CI099  •  1RV24CI096",
            0.5, 5.7, 9.5, 0.4, font_size=Pt(13), color=LIGHT_GREY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Agenda
# ════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
content_slide_header(s2, "Agenda")

items = [
    ("01", "Introduction & Motivation"),
    ("02", "Problem Statement"),
    ("03", "Algorithm 1 — Branch and Bound"),
    ("04", "Algorithm 2 — Held-Karp (Dynamic Programming)"),
    ("05", "Algorithm 3 — Christofides Approximation"),
    ("06", "Experimental Results & Comparison"),
    ("07", "Conclusion & Future Work"),
]
col_w, col_h = 5.7, 0.7
for i, (num, label) in enumerate(items):
    row, col = divmod(i, 2)
    lx = 0.5 + col * 6.4
    ty = 1.25 + row * 0.88
    add_rect(s2, lx, ty, col_w, col_h, fill_color=CARD_BG, line_color=ACCENT1, line_width=Pt(0.8))
    add_textbox(s2, num,   lx+0.12, ty+0.08, 0.7, col_h-0.15, font_size=Pt(22), bold=True, color=ACCENT2)
    add_textbox(s2, label, lx+0.82, ty+0.15, col_w-0.95, col_h-0.2, font_size=Pt(14), color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Introduction
# ════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
content_slide_header(s3, "Introduction & Motivation")

add_textbox(s3, "Why Route Optimisation Matters", 0.4, 1.15, 12.5, 0.45,
            font_size=Pt(17), bold=True, color=ACCENT2)

intro_bullets = [
    "Modern logistics, delivery networks, and airline scheduling all depend on finding efficient routes.",
    "The Travelling Salesman Problem (TSP) is a cornerstone of combinatorial optimisation.",
    "Real-world applications: e-commerce delivery, ride-sharing, supply-chain management, GPS navigation.",
    "As city counts grow, naive solutions become computationally infeasible — smart algorithms are essential.",
    "This project implements and benchmarks three major algorithmic paradigms to tackle TSP effectively.",
]
add_bullet_card(s3, 0.4, 1.65, 12.5, 5.5, "Key Motivations", intro_bullets)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Problem Statement
# ════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
content_slide_header(s4, "Problem Statement")

add_textbox(s4,
    "Given N cities and a distance matrix D, find the minimum-cost Hamiltonian cycle that "
    "visits every city exactly once and returns to the origin.",
    0.4, 1.15, 12.5, 0.9, font_size=Pt(16), color=WHITE)

cards = [
    ("Brute Force Limitation", ["O(N!) time complexity", "Infeasible beyond ~12 cities", "No pruning"]),
    ("Exact Algorithms",       ["Optimal solutions guaranteed", "Exponential worst-case", "Feasible for small N"]),
    ("Approximation Methods",  ["Near-optimal in polynomial time", "Sacrifices some accuracy", "Scales to large N"]),
    ("Our Goal",               ["Implement all three paradigms", "Benchmark on same datasets", "Quantify trade-offs"]),
]
positions = [(0.4,2.3),(4.7,2.3),(8.9,2.3),(0.4,5.1)]
widths    = [4.0,4.0,4.0,12.5]
for (lx,ty), w, (title, bullets) in zip(positions, widths, cards):
    add_bullet_card(s4, lx, ty, w, 2.4, title, bullets)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Branch and Bound
# ════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
content_slide_header(s5, "Algorithm 1 — Branch and Bound")

add_bullet_card(s5, 0.4, 1.15, 5.9, 5.95, "How it Works", [
    "Systematically explores all possible tours as a search tree.",
    "Each node = a partial tour; each branch = adding the next city.",
    "A lower-bound is calculated at every node using min-edge costs.",
    "Branches whose bound ≥ current best cost are PRUNED.",
    "Guarantees the globally optimal (exact) solution.",
    "Implementation uses recursive DFS with a nonlocal best tracker.",
])

add_bullet_card(s5, 6.6, 1.15, 6.3, 2.8, "Complexity", [
    "Time:  O(N!)  worst-case",
    "Space: O(N)   recursion stack",
    "Pruning significantly reduces average-case exploration.",
    "Practical limit: ~12–14 cities before runtime spikes.",
])

add_bullet_card(s5, 6.6, 4.2, 6.3, 2.9, "Results Snapshot", [
    "5 nodes  → 0.00012 s  | Cost 136.21",
    "8 nodes  → 0.00477 s  | Cost 287.54",
    "10 nodes → 0.06930 s  | Cost 244.91",
    "12 nodes → 3.59814 s  | Cost 328.56",
    "14 nodes → SKIPPED (too slow)",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Held-Karp
# ════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
content_slide_header(s6, "Algorithm 2 — Held-Karp (Dynamic Programming)")

add_bullet_card(s6, 0.4, 1.15, 5.9, 5.95, "How it Works", [
    "Uses bitmask DP to represent subsets of visited cities.",
    "State: dp[node][visited_mask] = min cost to reach 'node' after visiting 'mask'.",
    "Memoisation avoids recomputing overlapping sub-problems.",
    "After filling all states, back-tracks to reconstruct the optimal path.",
    "Delivers the exact optimal tour — identical result to B&B.",
    "Much faster than B&B thanks to overlapping-subproblem elimination.",
])

add_bullet_card(s6, 6.6, 1.15, 6.3, 2.8, "Complexity", [
    "Time:  O(2^N · N²)",
    "Space: O(2^N · N)  for the memo table",
    "Exponential, but far better than O(N!) in practice.",
    "Practical limit: ~20–22 cities (memory constraint).",
])

add_bullet_card(s6, 6.6, 4.2, 6.3, 2.9, "Results Snapshot", [
    "5 nodes  → 0.00006 s  | Cost 136.21",
    "8 nodes  → 0.00190 s  | Cost 287.54",
    "10 nodes → 0.00820 s  | Cost 244.91",
    "12 nodes → 0.06358 s  | Cost 328.56",
    "14 nodes → 0.31274 s  | Cost 395.88",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Christofides
# ════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
content_slide_header(s7, "Algorithm 3 — Christofides Approximation")

add_bullet_card(s7, 0.4, 1.15, 5.9, 5.95, "How it Works", [
    "Step 1: Build a Minimum Spanning Tree (MST) of the graph.",
    "Step 2: Find all odd-degree vertices in the MST.",
    "Step 3: Compute min-weight perfect matching on odd vertices.",
    "Step 4: Combine MST + matching → Eulerian multigraph.",
    "Step 5: Find an Eulerian circuit in the multigraph.",
    "Step 6: Shortcut repeated vertices → Hamiltonian cycle.",
    "Guarantees tour cost ≤ 1.5 × optimal (triangle inequality).",
])

add_bullet_card(s7, 6.6, 1.15, 6.3, 2.8, "Complexity", [
    "Time:  O(N³)  — polynomial!",
    "Space: O(N²)",
    "Approximation ratio: ≤ 1.5× optimal",
    "Requires triangle inequality to hold (Euclidean graphs).",
])

add_bullet_card(s7, 6.6, 4.2, 6.3, 2.9, "Results Snapshot", [
    "5 nodes  → 0.00304 s  | Error  3.24%",
    "8 nodes  → 0.00107 s  | Error  2.79%",
    "10 nodes → 0.00201 s  | Error  5.68%",
    "12 nodes → 0.00175 s  | Error  0.00%",
    "14 nodes → 0.00178 s  | Error 10.94%",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Comparison Table
# ════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
content_slide_header(s8, "Experimental Results — Comparative Analysis")

# Table data
headers = ["Nodes", "B&B Time (s)", "HK Time (s)", "Christofides Time (s)", "Optimal Cost", "CF Error"]
rows = [
    ["5",  "0.00012", "0.00006", "0.00304", "136.21", "3.24%"],
    ["8",  "0.00477", "0.00190", "0.00107", "287.54", "2.79%"],
    ["10", "0.06930", "0.00820", "0.00201", "244.91", "5.68%"],
    ["12", "3.59814", "0.06358", "0.00175", "328.56", "0.00%"],
    ["14", "Skipped", "0.31274", "0.00178", "395.88","10.94%"],
]

col_widths = [1.0, 1.9, 1.8, 2.5, 1.9, 1.6]
col_starts = [0.4]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1]+w)

row_h = 0.52
header_top = 1.15

for ci, (h, lx, w) in enumerate(zip(headers, col_starts, col_widths)):
    add_rect(s8, lx, header_top, w-0.05, row_h, fill_color=RGBColor(0x00,0x4A,0x6E), line_color=ACCENT1, line_width=Pt(0.5))
    add_textbox(s8, h, lx+0.05, header_top+0.08, w-0.1, row_h-0.15,
                font_size=Pt(11), bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    ty = header_top + (ri+1)*row_h
    bg = CARD_BG if ri % 2 == 0 else DARK_BG
    for ci, (val, lx, w) in enumerate(zip(row, col_starts, col_widths)):
        add_rect(s8, lx, ty, w-0.05, row_h, fill_color=bg, line_color=RGBColor(0x22,0x44,0x66), line_width=Pt(0.4))
        fc = ACCENT2 if val == "Skipped" else WHITE
        add_textbox(s8, val, lx+0.05, ty+0.09, w-0.1, row_h-0.15,
                    font_size=Pt(12), color=fc, align=PP_ALIGN.CENTER)

# Key takeaway
add_textbox(s8, "★  Held-Karp is fastest exact algorithm.  ★  Christofides is fastest overall — within ~11% of optimal.",
            0.4, 4.95, 12.5, 0.55, font_size=Pt(13), bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Algorithm Comparison Summary
# ════════════════════════════════════════════════════════════════════════════
s9 = add_slide()
content_slide_header(s9, "Algorithm Comparison Summary")

criteria = ["Optimality", "Time Complexity", "Space Complexity", "Scalability", "Best For"]
bb   = ["Exact (optimal)", "O(N!) worst-case", "O(N) stack", "Low (≤12 cities)", "Small exact problems"]
hk   = ["Exact (optimal)", "O(2^N · N²)",      "O(2^N · N)",  "Medium (≤20 cities)", "Moderate-size exact"]
cf   = ["≤1.5× optimal",   "O(N³)",             "O(N²)",       "High (any size)", "Large-scale routing"]

col_w = 3.8
for ci, (algo, data, clr) in enumerate([("Branch & Bound", bb, ACCENT2),
                                         ("Held-Karp",      hk, ACCENT1),
                                         ("Christofides",   cf, hex_rgb("#00E5A0"))]):
    lx = 0.4 + ci * 4.25
    add_rect(s9, lx, 1.15, col_w, 0.48, fill_color=clr)
    add_textbox(s9, algo, lx, 1.18, col_w, 0.42,
                font_size=Pt(17), bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    for ri, (crit, val) in enumerate(zip(criteria, data)):
        ty = 1.72 + ri * 1.0
        add_rect(s9, lx, ty, col_w, 0.42, fill_color=CARD_BG, line_color=clr, line_width=Pt(0.6))
        add_textbox(s9, crit, lx+0.1, ty+0.03, col_w-0.2, 0.22, font_size=Pt(10), bold=True, color=LIGHT_GREY)
        add_textbox(s9, val,  lx+0.1, ty+0.22, col_w-0.2, 0.22, font_size=Pt(11), color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Conclusion
# ════════════════════════════════════════════════════════════════════════════
s10 = add_slide()
content_slide_header(s10, "Conclusion & Future Work")

add_bullet_card(s10, 0.4, 1.15, 6.2, 5.95, "Key Conclusions", [
    "All three algorithms successfully solve TSP on Euclidean graphs.",
    "Branch & Bound guarantees optimality but is limited to small N.",
    "Held-Karp is the best exact algorithm — exponential but tractable up to ~20 cities.",
    "Christofides delivers near-optimal solutions almost instantaneously.",
    "Euclidean distances satisfy the triangle inequality, making Christofides ≤1.5× optimal.",
    "For real-world large-scale routing, Christofides is the most practical choice.",
])

add_bullet_card(s10, 6.9, 1.15, 6.0, 5.95, "Future Work", [
    "Implement Genetic Algorithms and Ant Colony Optimisation for comparison.",
    "Extend to asymmetric TSP (non-Euclidean cost matrices).",
    "Integrate with real map APIs (Google Maps, OpenStreetMap).",
    "Parallelise Branch & Bound using multi-threading.",
    "Build an interactive visualisation dashboard for tour animation.",
    "Apply algorithms to vehicle routing (multi-depot, time windows).",
])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Thank You
# ════════════════════════════════════════════════════════════════════════════
s11 = add_slide()
fill_bg(s11)
add_gradient_bar(s11, 0, 0,    13.33, 0.12, ACCENT1)
add_gradient_bar(s11, 0, 7.38, 13.33, 0.12, ACCENT2)

# Decorative circle
c2 = s11.shapes.add_shape(9, Inches(9.5), Inches(2.5), Inches(3.5), Inches(3.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x00,0x3A,0x5E)
c2.line.fill.background()

add_textbox(s11, "THANK YOU", 1.0, 2.2, 8.5, 1.3,
            font_size=Pt(56), bold=True, color=ACCENT1, align=PP_ALIGN.LEFT)
add_rect(s11, 1.0, 3.7, 5.5, 0.07, fill_color=ACCENT2)
add_textbox(s11, "Rhea Parthiban  •  Rishi Agarwal  •  Rangappagari John Niranjan",
            1.0, 3.9, 8.5, 0.55, font_size=Pt(16), bold=True, color=ACCENT2)
add_textbox(s11, "Design and Analysis of Algorithms  |  CD343AI  |  Dept. of AIML",
            1.0, 4.5, 8.5, 0.5, font_size=Pt(14), color=LIGHT_GREY)
add_textbox(s11, "Questions & Discussion Welcome", 1.0, 5.5, 8.5, 0.5,
            font_size=Pt(15), italic=True, color=LIGHT_GREY)

# ────────────────────────────────────────────────────────────────────────────
OUTPUT = r"c:\Users\rishi\OneDrive\Desktop\antigravity\TSP_DAA_Presentation.pptx"
prs.save(OUTPUT)
print(f"DONE. Saved: {OUTPUT}")
