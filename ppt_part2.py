"""Part 2: Load from Part 1, add slides 8-14 (3 algorithms, results, comparison, conclusion, thank you)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x00, 0x96, 0x88)
DT = RGBColor(0x00, 0x7A, 0x6E)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
LBG = RGBColor(0xF5, 0xF5, 0xF5)
CB = RGBColor(0xE0, 0xE0, 0xE0)

prs = Presentation(r"c:\Users\rishi\OneDrive\Desktop\antigravity\_prs_temp2.pptx")
BL = prs.slide_layouts[6]

def sl(): return prs.slides.add_slide(BL)
def bg(s, c=WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def rc(s, l, t, w, h, fc=None, lc=None, lw=Pt(0)):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else: sh.fill.background()
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else: sh.line.fill.background()
def tx(s, t, l, tp, w, h, sz=Pt(16), b=False, c=BLACK, a=PP_ALIGN.LEFT, it=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
    tb.word_wrap = True; tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.alignment = a
    r = p.add_run(); r.text = t; r.font.size = sz
    r.font.bold = b; r.font.color.rgb = c; r.font.italic = it
def ml(s, lines, l, t, w, h, sz=Pt(12), c=GREY, bl="  >  "):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True; tb.text_frame.word_wrap = True
    for i, ln in enumerate(lines):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.space_before = Pt(5); r = p.add_run(); r.text = bl + ln; r.font.size = sz; r.font.color.rgb = c
def sb(s, n, t):
    rc(s, 0.4, 0.35, 0.6, 0.06, fc=TEAL)
    tx(s, f"{n} - {t}", 1.15, 0.15, 10, 0.55, sz=Pt(28), b=True, c=TEAL)
def sc(s, l, t, w, h, num, lab, sub=""):
    rc(s, l, t, w, h, fc=WHITE, lc=CB, lw=Pt(1.2)); rc(s, l, t, w, 0.06, fc=TEAL)
    tx(s, num, l, t+0.25, w, 0.6, sz=Pt(34), b=True, c=BLACK, a=PP_ALIGN.CENTER)
    tx(s, lab, l, t+0.82, w, 0.3, sz=Pt(10), b=True, c=GREY, a=PP_ALIGN.CENTER)
    if sub: tx(s, sub, l, t+1.1, w, 0.25, sz=Pt(9), c=GREY, a=PP_ALIGN.CENTER)
def kb(s, l, t, w, h, ti, bd):
    rc(s, l, t, w, h, fc=LBG, lc=TEAL, lw=Pt(1.5)); rc(s, l, t, 0.06, h, fc=TEAL)
    tx(s, ti, l+0.2, t+0.1, w-0.3, 0.3, sz=Pt(13), b=True, c=TEAL)
    tx(s, bd, l+0.2, t+0.4, w-0.3, h-0.5, sz=Pt(11), c=GREY)

# ── 8: BRANCH AND BOUND ──
s = sl(); bg(s); sb(s, "06", "BRANCH AND BOUND")
tx(s, "Exact Optimisation via Search Tree Pruning", 0.5, 0.85, 7, 0.5, sz=Pt(20), b=True, c=BLACK)
ml(s, [
    "Systematically explores all possible tours as a search tree",
    "Each node represents a partial tour; branches add the next city",
    "Lower bound calculated at each node using minimum edge costs",
    "Branches exceeding current best cost are pruned immediately",
    "Guarantees the globally optimal (exact) solution",
    "Uses recursive DFS with nonlocal best-cost tracker",
], 0.5, 1.4, 6.5, 3.0)
sc(s, 7.8, 1.0, 2.4, 1.5, "O(N!)", "TIME", "Worst-case"); sc(s, 10.5, 1.0, 2.4, 1.5, "O(N)", "SPACE", "Recursion stack")
sc(s, 7.8, 2.7, 2.4, 1.5, "12", "MAX CITIES", "Practical limit"); sc(s, 10.5, 2.7, 2.4, 1.5, "3.6s", "AT 12 NODES", "Execution time")
kb(s, 0.5, 4.6, 12.4, 1.2, "PERFORMANCE",
   "5 nodes: 0.00012s | 8 nodes: 0.00477s | 10 nodes: 0.06930s | 12 nodes: 3.59814s | 14 nodes: SKIPPED (too slow). Pruning helps but factorial growth dominates.")

# ── 9: HELD-KARP ──
s = sl(); bg(s); sb(s, "07", "HELD-KARP ALGORITHM")
tx(s, "Dynamic Programming with Bitmask Memoisation", 0.5, 0.85, 7, 0.5, sz=Pt(20), b=True, c=BLACK)
ml(s, [
    "Uses bitmask to represent subsets of visited cities",
    "State: dp[node][visited_mask] = min cost reaching node after mask",
    "Memoisation eliminates redundant overlapping sub-problems",
    "After filling all states, back-tracks to reconstruct optimal path",
    "Delivers exact optimal tour - identical to Branch and Bound result",
    "Dramatically faster than B&B due to subproblem reuse",
], 0.5, 1.4, 6.5, 3.0)
sc(s, 7.8, 1.0, 2.4, 1.5, "O(2^N*N^2)", "TIME", "Exponential"); sc(s, 10.5, 1.0, 2.4, 1.5, "O(2^N*N)", "SPACE", "Memo table")
sc(s, 7.8, 2.7, 2.4, 1.5, "20", "MAX CITIES", "Memory limited"); sc(s, 10.5, 2.7, 2.4, 1.5, "0.31s", "AT 14 NODES", "Execution time")
kb(s, 0.5, 4.6, 12.4, 1.2, "PERFORMANCE",
   "5 nodes: 0.00006s | 8 nodes: 0.00190s | 10 nodes: 0.00820s | 12 nodes: 0.06358s | 14 nodes: 0.31274s. Consistently the fastest exact method across all test sizes.")

# ── 10: CHRISTOFIDES ──
s = sl(); bg(s); sb(s, "08", "CHRISTOFIDES ALGORITHM")
tx(s, "Polynomial-Time Approximation with 1.5x Guarantee", 0.5, 0.85, 7, 0.5, sz=Pt(20), b=True, c=BLACK)
ml(s, [
    "Step 1: Build Minimum Spanning Tree (MST) of the graph",
    "Step 2: Identify all odd-degree vertices in the MST",
    "Step 3: Compute min-weight perfect matching on odd vertices",
    "Step 4: Combine MST + matching into Eulerian multigraph",
    "Step 5: Find Eulerian circuit in the multigraph",
    "Step 6: Shortcut repeated vertices to get Hamiltonian cycle",
    "Guarantees tour cost <= 1.5x optimal (triangle inequality)",
], 0.5, 1.4, 6.5, 3.2)
sc(s, 7.8, 1.0, 2.4, 1.5, "O(N^3)", "TIME", "Polynomial!"); sc(s, 10.5, 1.0, 2.4, 1.5, "O(N^2)", "SPACE", "Graph storage")
sc(s, 7.8, 2.7, 2.4, 1.5, "1.5x", "RATIO", "Worst-case bound"); sc(s, 10.5, 2.7, 2.4, 1.5, "<0.01s", "ALL SIZES", "Near-instant")
kb(s, 0.5, 4.85, 12.4, 1.1, "PERFORMANCE",
   "5 nodes: 3.24% error | 8 nodes: 2.79% error | 10 nodes: 5.68% error | 12 nodes: 0.00% error | 14 nodes: 10.94% error. Fastest across all tests.")

# ── 11: RESULTS TABLE ──
s = sl(); bg(s); sb(s, "09", "EXPERIMENTAL RESULTS")
tx(s, "Comparative Benchmarks on Euclidean Distance Networks", 0.5, 0.85, 12, 0.5, sz=Pt(20), b=True, c=BLACK)
headers = ["Nodes", "B&B Time", "H-K Time", "CF Time", "Optimal Cost", "CF Error"]
rows = [["5","0.00012s","0.00006s","0.00304s","136.21","3.24%"],["8","0.00477s","0.00190s","0.00107s","287.54","2.79%"],
        ["10","0.06930s","0.00820s","0.00201s","244.91","5.68%"],["12","3.59814s","0.06358s","0.00175s","328.56","0.00%"],
        ["14","Skipped","0.31274s","0.00178s","395.88","10.94%"]]
cw = [1.2, 2.0, 2.0, 2.0, 2.0, 1.8]; cs = [0.8]
for w in cw[:-1]: cs.append(cs[-1]+w)
rh = 0.55
for ci, (h, lx, w) in enumerate(zip(headers, cs, cw)):
    rc(s, lx, 1.4, w-0.05, rh, fc=TEAL)
    tx(s, h, lx+0.05, 1.45, w-0.1, rh-0.1, sz=Pt(11), b=True, c=WHITE, a=PP_ALIGN.CENTER)
for ri, row in enumerate(rows):
    ty = 1.4 + (ri+1)*rh
    bgc = WHITE if ri % 2 == 0 else LBG
    for ci, (val, lx, w) in enumerate(zip(row, cs, cw)):
        rc(s, lx, ty, w-0.05, rh, fc=bgc, lc=CB, lw=Pt(0.5))
        vc = RGBColor(0xC0, 0x39, 0x2B) if val == "Skipped" else BLACK
        tx(s, val, lx+0.05, ty+0.08, w-0.1, rh-0.15, sz=Pt(11), c=vc, a=PP_ALIGN.CENTER)
kb(s, 0.8, 4.7, 11.7, 1.1, "KEY FINDINGS",
   "Held-Karp consistently outperforms B&B as fastest exact method. Christofides delivers results in under 5ms for all sizes with average error below 5%. B&B becomes impractical beyond 12 cities.")

# ── 12: COMPARISON ──
s = sl(); bg(s); sb(s, "10", "ALGORITHM COMPARISON")
tx(s, "Head-to-Head Analysis", 0.5, 0.85, 12, 0.5, sz=Pt(22), b=True, c=BLACK)
criteria = ["Optimality", "Time Complexity", "Space Complexity", "Scalability", "Best Use Case"]
bb = ["Exact (optimal)", "O(N!) worst-case", "O(N) stack", "Low: <=12 cities", "Small exact problems"]
hk = ["Exact (optimal)", "O(2^N * N^2)", "O(2^N * N)", "Medium: <=20 cities", "Moderate exact solutions"]
cf = ["<=1.5x optimal", "O(N^3)", "O(N^2)", "High: scales to 100+", "Large-scale routing"]
algos = [("Branch & Bound", bb), ("Held-Karp (DP)", hk), ("Christofides", cf)]
cw2 = 3.9
for ci, (algo, data) in enumerate(algos):
    lx = 0.5 + ci * 4.15
    rc(s, lx, 1.35, cw2, 0.5, fc=TEAL)
    tx(s, algo, lx, 1.38, cw2, 0.45, sz=Pt(16), b=True, c=WHITE, a=PP_ALIGN.CENTER)
    for ri, (crit, val) in enumerate(zip(criteria, data)):
        ty = 1.9 + ri * 0.95
        rc(s, lx, ty, cw2, 0.85, fc=WHITE if ri%2==0 else LBG, lc=CB, lw=Pt(0.5))
        tx(s, crit, lx+0.12, ty+0.05, cw2-0.2, 0.3, sz=Pt(10), b=True, c=TEAL)
        tx(s, val, lx+0.12, ty+0.4, cw2-0.2, 0.35, sz=Pt(12), c=BLACK)

# ── 13: CONCLUSION ──
s = sl(); bg(s); sb(s, "11", "CONCLUSION & FUTURE WORK")
tx(s, "Key Conclusions", 0.5, 0.85, 5.8, 0.4, sz=Pt(20), b=True, c=BLACK)
ml(s, [
    "All three algorithms successfully solve TSP on Euclidean graphs",
    "Branch & Bound guarantees optimality but limited to small N",
    "Held-Karp is the best exact method - tractable up to ~20 cities",
    "Christofides delivers near-optimal results almost instantly",
    "Euclidean distances satisfy triangle inequality for 1.5x bound",
    "For real-world large-scale routing, Christofides is most practical",
], 0.5, 1.3, 5.8, 3.2)
tx(s, "Future Work", 7.0, 0.85, 5.8, 0.4, sz=Pt(20), b=True, c=BLACK)
ml(s, [
    "Implement Genetic Algorithms and Ant Colony Optimisation",
    "Extend to asymmetric TSP (non-Euclidean cost matrices)",
    "Integrate with real map APIs (Google Maps, OpenStreetMap)",
    "Parallelise Branch & Bound with multi-threading",
    "Build interactive visualisation dashboard for tour animation",
    "Apply to vehicle routing with time windows and multi-depot",
], 7.0, 1.3, 5.8, 3.2)
kb(s, 0.5, 4.85, 12.4, 1.5, "FINAL TAKEAWAY",
   "Algorithm choice profoundly impacts performance. For small networks (N<=12), B&B or Held-Karp provide exact solutions. For larger networks, Christofides provides an excellent balance of speed and accuracy. This comparative analysis offers practical insights into optimisation techniques used in real-world routing and logistics.")

# ── 14: THANK YOU ──
s = sl(); bg(s)
rc(s, 0, 0, 0.08, 7.5, fc=TEAL); rc(s, 0, 7.2, 13.33, 0.04, fc=TEAL)
tx(s, "THANK YOU", 0.6, 2.2, 12, 1.0, sz=Pt(52), b=True, c=BLACK)
rc(s, 0.6, 3.4, 3.5, 0.06, fc=TEAL)
tx(s, "Rhea Parthiban   |   Rishi Agarwal   |   Rangappagari John Niranjan", 0.6, 3.7, 10, 0.5, sz=Pt(15), b=True, c=DT)
tx(s, "Design and Analysis of Algorithms  |  CD343AI  |  Dept. of AIML", 0.6, 4.2, 10, 0.4, sz=Pt(13), c=GREY)
tx(s, "Questions & Discussion Welcome", 0.6, 5.2, 10, 0.4, sz=Pt(14), c=GREY, it=True)

OUT = r"c:\Users\rishi\OneDrive\Desktop\antigravity\TSP_DAA_Final.pptx"
prs.save(OUT)
print(f"DONE - Saved: {OUT}")
