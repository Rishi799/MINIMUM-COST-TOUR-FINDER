"""Part 1: Slides 1-7 (Title, Agenda, Intro, Problem, Objectives, Lit Review, Methodology)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x00, 0x96, 0x88)
DT = RGBColor(0x00, 0x7A, 0x6E)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
LBG = RGBColor(0xF5, 0xF5, 0xF5)
CB = RGBColor(0xE0, 0xE0, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]

def sl(): return prs.slides.add_slide(BL)
def bg(s, c=WHITE):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c
def rc(s, l, t, w, h, fc=None, lc=None, lw=Pt(0)):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else: sh.fill.background()
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else: sh.line.fill.background()
def tx(s, t, l, tp, w, h, sz=Pt(16), b=False, c=BLACK, a=PP_ALIGN.LEFT, it=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
    tb.word_wrap = True; tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.alignment = a
    r = p.add_run(); r.text = t; r.font.size = sz
    r.font.bold = b; r.font.color.rgb = c; r.font.italic = it
def ml(s, lines, l, t, w, h, sz=Pt(12), c=GREY, bl="  >  "):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True; tb.text_frame.word_wrap = True
    for i, ln in enumerate(lines):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.space_before = Pt(5); r = p.add_run(); r.text = bl + ln; r.font.size = sz; r.font.color.rgb = c
def sb(s, n, t):
    rc(s, 0.4, 0.35, 0.6, 0.06, fc=TEAL)
    tx(s, f"{n} - {t}", 1.15, 0.15, 10, 0.55, sz=Pt(28), b=True, c=TEAL)
def sc(s, l, t, w, h, num, lab, sub=""):
    rc(s, l, t, w, h, fc=WHITE, lc=CB, lw=Pt(1.2)); rc(s, l, t, w, 0.06, fc=TEAL)
    tx(s, num, l, t+0.25, w, 0.6, sz=Pt(34), b=True, c=BLACK, a=PP_ALIGN.CENTER)
    tx(s, lab, l, t+0.82, w, 0.3, sz=Pt(10), b=True, c=GREY, a=PP_ALIGN.CENTER)
    if sub: tx(s, sub, l, t+1.1, w, 0.25, sz=Pt(9), c=GREY, a=PP_ALIGN.CENTER)
def kb(s, l, t, w, h, ti, bd):
    rc(s, l, t, w, h, fc=LBG, lc=TEAL, lw=Pt(1.5)); rc(s, l, t, 0.06, h, fc=TEAL)
    tx(s, ti, l+0.2, t+0.1, w-0.3, 0.3, sz=Pt(13), b=True, c=TEAL)
    tx(s, bd, l+0.2, t+0.4, w-0.3, h-0.5, sz=Pt(11), c=GREY)

# ── 1: TITLE ──
s = sl(); bg(s)
rc(s, 0, 0, 0.08, 7.5, fc=TEAL); rc(s, 0, 7.2, 13.33, 0.04, fc=TEAL)
tx(s, "MINIMUM COST TOUR FINDER", 0.6, 1.6, 12, 0.8, sz=Pt(42), b=True, c=BLACK)
tx(s, "For Multi-City Travel Networks", 0.6, 2.45, 12, 0.6, sz=Pt(24), c=TEAL, b=True)
rc(s, 0.6, 3.2, 4, 0.05, fc=TEAL)
tx(s, "Design and Analysis of Algorithms  |  CD343AI", 0.6, 3.5, 10, 0.4, sz=Pt(14), c=GREY)
tx(s, "Department of Artificial Intelligence and Machine Learning", 0.6, 3.95, 10, 0.4, sz=Pt(12), c=GREY, it=True)
tx(s, "Rhea Parthiban   |   Rishi Agarwal   |   Rangappagari John Niranjan", 0.6, 5.0, 10, 0.4, sz=Pt(15), b=True, c=DT)
tx(s, "1RV24CI097   |   1RV24CI099   |   1RV24CI096", 0.6, 5.45, 10, 0.4, sz=Pt(12), c=GREY)

# ── 2: AGENDA ──
s = sl(); bg(s); sb(s, "00", "AGENDA")
items = ["Introduction & Relevance", "Problem Statement", "Objectives",
         "Literature Review", "Proposed Methodology",
         "Algorithm Deep-Dives (B&B, Held-Karp, Christofides)",
         "Experimental Results & Comparison",
         "Conclusion & Future Work"]
for i, item in enumerate(items):
    row, col = divmod(i, 2)
    lx, ty = 0.5 + col*6.3, 1.1 + row*1.4
    rc(s, lx, ty, 5.8, 1.1, fc=WHITE, lc=CB, lw=Pt(1)); rc(s, lx, ty, 0.06, 1.1, fc=TEAL)
    tx(s, f"0{i+1}", lx+0.2, ty+0.15, 0.6, 0.5, sz=Pt(26), b=True, c=TEAL)
    tx(s, item, lx+0.85, ty+0.3, 4.8, 0.5, sz=Pt(14), c=BLACK)

# ── 3: INTRODUCTION ──
s = sl(); bg(s); sb(s, "01", "INTRODUCTION")
tx(s, "The Route Optimisation Challenge", 0.5, 0.85, 7, 0.5, sz=Pt(22), b=True, c=BLACK)
ml(s, [
    "Modern logistics and transportation systems require efficient route planning across complex multi-city networks.",
    "The Travelling Salesman Problem (TSP) is a fundamental combinatorial optimisation problem: find the shortest tour visiting every city exactly once and returning to the start.",
    "Applications span delivery networks, airline scheduling, ride-sharing platforms, supply-chain management, and AI-based navigation.",
    "As city counts grow, naive solutions become computationally infeasible - smart algorithms are essential.",
], 0.5, 1.4, 6.5, 3.0)
sc(s, 7.8, 1.0, 2.4, 1.5, "3", "ALGORITHMS", "Implemented"); sc(s, 10.5, 1.0, 2.4, 1.5, "5", "TEST SIZES", "5 to 14 cities")
sc(s, 7.8, 2.7, 2.4, 1.5, "N!", "BRUTE FORCE", "Factorial growth"); sc(s, 10.5, 2.7, 2.4, 1.5, "<1s", "APPROX TIME", "Christofides")
kb(s, 0.5, 4.6, 12.4, 1.2, "PROJECT SCOPE",
   "Comparative implementation and benchmarking of three TSP algorithms - Branch and Bound, Held-Karp, and Christofides - covering exact and approximation paradigms to analyse trade-offs between optimality, speed, and scalability.")

# ── 4: PROBLEM STATEMENT ──
s = sl(); bg(s); sb(s, "02", "PROBLEM STATEMENT")
tx(s, "Defining the Challenge", 0.5, 0.85, 12, 0.5, sz=Pt(22), b=True, c=BLACK)
tx(s, "Given N cities and a pairwise distance matrix D, find the minimum-cost Hamiltonian cycle visiting every city exactly once and returning to the origin.", 0.5, 1.35, 12, 0.7, sz=Pt(14), c=GREY)
lims = [
    ("Brute Force Limitation", "O(N!) time complexity makes it infeasible beyond ~10 cities. No pruning or intelligence applied."),
    ("Exact Algorithm Trade-off", "Guarantee optimal solutions but suffer exponential worst-case time as the number of cities grows."),
    ("Approximation Trade-off", "Trade optimality for speed. Near-optimal solutions in polynomial time but sacrifice some precision."),
    ("The Core Challenge", "Evaluate and compare different algorithmic strategies to find solutions that are optimal or near-optimal within a reasonable time."),
]
for i, (t, d) in enumerate(lims):
    row, col = divmod(i, 2)
    lx, ty = 0.5 + col*6.4, 2.2 + row*2.3
    rc(s, lx, ty, 6.1, 2.0, fc=WHITE, lc=CB, lw=Pt(1)); rc(s, lx, ty, 6.1, 0.06, fc=TEAL)
    tx(s, t, lx+0.15, ty+0.15, 5.8, 0.35, sz=Pt(14), b=True, c=TEAL)
    tx(s, d, lx+0.15, ty+0.55, 5.8, 1.2, sz=Pt(12), c=GREY)

# ── 5: OBJECTIVES ──
s = sl(); bg(s); sb(s, "03", "OBJECTIVES")
tx(s, "What This Project Aims to Achieve", 0.5, 0.85, 12, 0.5, sz=Pt(22), b=True, c=BLACK)
objs = [
    ("Implement Branch and Bound", "Build a recursive DFS-based exact solver with lower-bound pruning for optimal tour computation."),
    ("Implement Held-Karp Algorithm", "Develop a bitmask dynamic programming solution with memoisation for exact TSP optimisation."),
    ("Implement Christofides Algorithm", "Create an approximation solver using MST, minimum matching, and Eulerian circuit techniques."),
    ("Benchmark and Compare", "Run all three algorithms on identical Euclidean datasets of varying sizes (5 to 14 cities)."),
    ("Analyse Trade-offs", "Evaluate execution time, solution accuracy, and scalability to identify the most suitable algorithm for different problem sizes."),
]
for i, (t, d) in enumerate(objs):
    ty = 1.35 + i*1.15
    rc(s, 0.5, ty, 12.4, 1.0, fc=WHITE, lc=CB, lw=Pt(1)); rc(s, 0.5, ty, 0.06, 1.0, fc=TEAL)
    tx(s, f"0{i+1}", 0.7, ty+0.08, 0.5, 0.4, sz=Pt(20), b=True, c=TEAL)
    tx(s, t, 1.3, ty+0.05, 5, 0.35, sz=Pt(14), b=True, c=BLACK)
    tx(s, d, 1.3, ty+0.42, 11.4, 0.5, sz=Pt(11), c=GREY)

# ── 6: LITERATURE REVIEW ──
s = sl(); bg(s); sb(s, "04", "LITERATURE REVIEW")
tx(s, "Key Research Papers and Findings", 0.5, 0.85, 12, 0.5, sz=Pt(22), b=True, c=BLACK)

# Table headers
headers = ["#", "Paper / Authors (Year)", "Method", "Key Findings"]
hw = [0.4, 4.5, 2.2, 5.2]
hs = [0.5]
for w in hw[:-1]: hs.append(hs[-1]+w)
rh = 0.5

for ci, (h, lx, w) in enumerate(zip(headers, hs, hw)):
    rc(s, lx, 1.35, w-0.04, rh, fc=TEAL)
    tx(s, h, lx+0.08, 1.4, w-0.15, rh-0.1, sz=Pt(11), b=True, c=WHITE, a=PP_ALIGN.LEFT)

papers = [
    ["1", "Applegate et al. (2006)\n\"The Traveling Salesman Problem:\nA Computational Study\"", "Branch & Cut\n+ LP Relaxation", "Solved instances up to 85,900 cities using\nbranch-and-cut. Established Concorde as the\ngold standard TSP solver."],
    ["2", "Bellman (1962), Held & Karp (1962)\n\"Dynamic Programming Treatment\nof the TSP\"", "Bitmask DP\nMemoisation", "Reduced TSP from O(N!) to O(2^N * N^2).\nFirst polynomial-space exact algorithm.\nFoundation of modern exact TSP solvers."],
    ["3", "Christofides (1976)\n\"Worst-Case Analysis of a New\nHeuristic for the TSP\"", "MST + Min\nMatching", "Proved 1.5x approximation ratio for metric\nTSP. Remained best approximation for 45\nyears until Karlin et al. (2021)."],
    ["4", "Dorigo & Gambardella (1997)\n\"Ant Colony System: A Cooperative\nLearning Approach to TSP\"", "Ant Colony\nOptimisation", "Bio-inspired meta-heuristic achieving near-\noptimal results on large instances. Showed\nswarm intelligence as viable TSP approach."],
]

for ri, row in enumerate(papers):
    ty = 1.35 + (ri+1)*rh*2.4
    bgc = WHITE if ri % 2 == 0 else LBG
    for ci, (val, lx, w) in enumerate(zip(row, hs, hw)):
        rc(s, lx, ty, w-0.04, rh*2.2, fc=bgc, lc=CB, lw=Pt(0.5))
        vc = TEAL if ci == 0 else BLACK
        tx(s, val, lx+0.08, ty+0.06, w-0.15, rh*2, sz=Pt(10), c=vc, b=(ci==0))

# ── 7: METHODOLOGY ──
s = sl(); bg(s); sb(s, "05", "PROPOSED METHODOLOGY")
tx(s, "Implementation Approach", 0.5, 0.85, 12, 0.5, sz=Pt(22), b=True, c=BLACK)

steps = [
    ("01", "Graph Representation", "Represent cities and inter-city distances as a weighted adjacency matrix. Generate Euclidean distance networks to satisfy the triangle inequality constraint."),
    ("02", "Algorithm Implementation", "Implement all three algorithms independently in Python: Branch and Bound (recursive DFS), Held-Karp (bitmask DP), and Christofides (NetworkX-based MST + matching)."),
    ("03", "Dataset Generation", "Create multiple test datasets with varying sizes (5, 8, 10, 12, 14 cities) using seeded random Euclidean point generation for reproducibility."),
    ("04", "Performance Measurement", "Measure execution time, tour cost, and solution quality (error % vs optimal) for each algorithm on every dataset."),
    ("05", "Comparative Analysis", "Compare algorithms across all metrics. Identify practical limits, scalability thresholds, and trade-offs between exact and approximate methods."),
    ("06", "Result Compilation", "Output optimal/near-optimal routes with total cost. Generate comparative tables and analysis for presentation."),
]
for i, (num, title, desc) in enumerate(steps):
    row, col = divmod(i, 2)
    lx, ty = 0.5 + col*6.4, 1.35 + row*1.9
    rc(s, lx, ty, 6.1, 1.7, fc=WHITE, lc=CB, lw=Pt(1)); rc(s, lx, ty, 6.1, 0.06, fc=TEAL)
    tx(s, num, lx+0.12, ty+0.15, 0.5, 0.4, sz=Pt(22), b=True, c=TEAL)
    tx(s, title, lx+0.65, ty+0.15, 5.2, 0.35, sz=Pt(14), b=True, c=BLACK)
    tx(s, desc, lx+0.15, ty+0.6, 5.8, 1.0, sz=Pt(10), c=GREY)

prs.save(r"c:\Users\rishi\OneDrive\Desktop\antigravity\_prs_temp2.pptx")
print("Part 1 done - 7 slides")
